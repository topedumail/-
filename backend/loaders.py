"""קריאת מסמכים מתיקיית הידע עם chunking שמכבד מבנה (כותרות, סעיפים)."""
from __future__ import annotations

import hashlib
import re
from dataclasses import dataclass, field
from pathlib import Path


@dataclass
class Document:
    source: str
    text: str
    title: str = ""  # אם המסמך מתחיל בכותרת ברורה


@dataclass
class Chunk:
    source: str            # שם הקובץ
    text: str              # טקסט ה-chunk
    section: str = ""      # כותרת הסעיף שאליו שייך ה-chunk (אם זוהה)
    position: int = 0      # מיקום ה-chunk במסמך (לסידור)


# ---------- קוראי קבצים ----------

def _read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _read_docx(path: Path) -> str:
    """קורא docx בצורה שמכבדת רמת כותרת — מוסיף Markdown # לפי Heading level."""
    from docx import Document as DocxDocument

    doc = DocxDocument(str(path))
    parts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style_name = (para.style.name or "").lower() if para.style else ""
        if "heading 1" in style_name or "title" in style_name:
            parts.append(f"# {text}")
        elif "heading 2" in style_name:
            parts.append(f"## {text}")
        elif "heading 3" in style_name:
            parts.append(f"### {text}")
        elif style_name.startswith("heading"):
            parts.append(f"#### {text}")
        else:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _read_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts: list[str] = []
    for page_num, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        text = text.strip()
        if text:
            parts.append(f"[עמוד {page_num}]\n{text}")
    return "\n\n".join(parts)


def _read_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        slide_parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        slide_parts.append(" | ".join(cells))
        if slide_parts:
            parts.append(f"## שקופית {slide_num}\n" + "\n".join(slide_parts))
    return "\n\n".join(parts)


_LOADERS = {
    ".txt": _read_text,
    ".md": _read_text,
    ".markdown": _read_text,
    ".docx": _read_docx,
    ".pdf": _read_pdf,
    ".pptx": _read_pptx,
}


# ---------- טעינת מסמכים ----------

def load_documents(folder: Path) -> list[Document]:
    """קורא את כל הקבצים הנתמכים בתיקייה ובתת-תיקיות, עם dedup לפי תוכן."""
    seen_hashes: dict[str, str] = {}
    docs: list[Document] = []
    for path in sorted(folder.rglob("*")):
        if not path.is_file():
            continue
        loader = _LOADERS.get(path.suffix.lower())
        if loader is None:
            continue
        try:
            text = loader(path).strip()
        except Exception as exc:
            print(f"[loaders] לא ניתן לקרוא את {path.name}: {exc}")
            continue
        if not text:
            continue
        rel = path.relative_to(folder).as_posix()
        content_hash = hashlib.sha256(text.encode("utf-8")).hexdigest()
        if content_hash in seen_hashes:
            print(f"[loaders] מדלג על כפילות: {rel} (זהה ל-{seen_hashes[content_hash]})")
            continue
        seen_hashes[content_hash] = rel
        # נסה לזהות כותרת ראשית
        title = _extract_title(text, fallback=Path(rel).stem)
        docs.append(Document(source=rel, text=text, title=title))
    return docs


def _extract_title(text: str, fallback: str) -> str:
    """מנסה לחלץ כותרת — שורה ראשונה אם היא קצרה, או heading של Markdown."""
    lines = text.split("\n", 5)
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("# "):
            return stripped.lstrip("# ").strip()
        if 3 <= len(stripped) <= 100 and not stripped.endswith("."):
            return stripped
        break
    return fallback


# ---------- Chunking שמכבד מבנה ----------

_HEADING_RE = re.compile(r"^(#{1,6})\s+(.+)$")
_SENTENCE_END_RE = re.compile(r"(?<=[.!?…])\s+(?=[א-ת0-9A-Z])")

CHUNK_TARGET = 600   # יעד גודל chunk בתווים
CHUNK_MAX = 900      # מקסימום קשיח
CHUNK_OVERLAP = 120  # הצפה בין chunks כדי לא לאבד הקשר בגבולות


def _split_by_sections(text: str) -> list[tuple[str, str]]:
    """מחלק טקסט לסעיפים לפי כותרות Markdown.

    Returns:
        list of (section_path, section_text). אם אין כותרות בכלל, מחזיר [("", text)].
    """
    sections: list[tuple[str, str]] = []
    current_section_parts: list[str] = []
    heading_stack: list[tuple[int, str]] = []  # (level, title)
    has_any_heading = False

    def flush_current():
        if current_section_parts:
            section_path = " › ".join(h[1] for h in heading_stack)
            section_text = "\n".join(current_section_parts).strip()
            if section_text:
                sections.append((section_path, section_text))

    for line in text.split("\n"):
        m = _HEADING_RE.match(line.strip())
        if m:
            has_any_heading = True
            flush_current()
            current_section_parts = []
            level = len(m.group(1))
            title = m.group(2).strip()
            # פוקע מ-stack כותרות באותה רמה או עמוקות יותר
            heading_stack = [h for h in heading_stack if h[0] < level]
            heading_stack.append((level, title))
        else:
            current_section_parts.append(line)
    flush_current()

    if not has_any_heading:
        return [("", text)]
    return sections


def _split_long_section(section_text: str) -> list[str]:
    """מפצל סעיף ארוך ל-chunks בגודל סביר, מנסה לפצל על גבולות משפט/פסקה."""
    if len(section_text) <= CHUNK_MAX:
        return [section_text]

    # פצל קודם לפי פסקאות
    paragraphs = [p.strip() for p in section_text.split("\n") if p.strip()]
    chunks: list[str] = []
    buf = ""

    for para in paragraphs:
        if len(buf) + len(para) + 1 <= CHUNK_TARGET:
            buf = f"{buf}\n{para}" if buf else para
            continue
        # ה-buf הנוכחי מספיק גדול — שמור והתחל חדש (עם overlap מהסוף)
        if buf:
            chunks.append(buf)
            # overlap — לקחת את סוף ה-buf כתחילה של ה-chunk הבא
            tail = buf[-CHUNK_OVERLAP:] if len(buf) > CHUNK_OVERLAP else ""
            buf = tail
        # אם הפסקה עצמה ארוכה מדי — פצל אותה לפי משפטים
        if len(para) > CHUNK_MAX:
            sentences = _SENTENCE_END_RE.split(para)
            for sent in sentences:
                if len(buf) + len(sent) + 1 <= CHUNK_TARGET:
                    buf = f"{buf} {sent}" if buf else sent
                else:
                    if buf:
                        chunks.append(buf)
                        tail = buf[-CHUNK_OVERLAP:] if len(buf) > CHUNK_OVERLAP else ""
                        buf = tail
                    # משפט בודד עדיין ארוך מדי — פצל ביד
                    if len(sent) > CHUNK_MAX:
                        for i in range(0, len(sent), CHUNK_TARGET):
                            chunks.append(sent[i : i + CHUNK_TARGET])
                        buf = ""
                    else:
                        buf = sent if not buf else f"{buf} {sent}"
        else:
            buf = para if not buf else f"{buf}\n{para}"

    if buf.strip():
        chunks.append(buf.strip())
    return [c for c in chunks if c.strip()]


def build_chunks(docs: list[Document]) -> list[Chunk]:
    """בונה רשימת Chunks עם metadata על הסעיף שכל chunk שייך אליו."""
    result: list[Chunk] = []
    for doc in docs:
        sections = _split_by_sections(doc.text)
        position = 0
        for section_path, section_text in sections:
            for chunk_text in _split_long_section(section_text):
                result.append(
                    Chunk(
                        source=doc.source,
                        text=chunk_text,
                        section=section_path,
                        position=position,
                    )
                )
                position += 1
    return result
